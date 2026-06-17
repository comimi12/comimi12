# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Emu, Pt
WD=r"C:\Users\owner\do-better-workspace-v2\00-inbox\voc-may-ppt"
SRC=r"C:\Users\owner\Desktop\교육팀\0. 온라인리뷰\★ 26년  VOC 리뷰★\온라인 리뷰_26.05월(05.01~05.31)_V2.pptx"
OUT=r"C:\Users\owner\Desktop\교육팀\0. 온라인리뷰\★ 26년  VOC 리뷰★\온라인 리뷰_26.05월(05.01~05.31)_V3.pptx"

prs=Presentation(SRC)
slide=prs.slides[1]

def setrun(shape,txt):
    p=shape.text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text=txt
        for r in p.runs[1:]: r.text=""
    else:
        p.text=txt

M={
 (1.35,1.33):"22",(2.00,1.33):"1,358",(3.00,1.33):"1,324",(4.00,1.33):"20",(4.82,1.33):"1.5%",
 (1.35,1.60):"38",(2.00,1.60):"1,394",(3.00,1.60):"1,343",(4.00,1.60):"30",(4.82,1.60):"2.2%",
 (1.35,1.87):"11",(2.00,1.87):"1,184",(3.00,1.87):"1,145",(4.00,1.87):"21",(4.82,1.87):"1.8%",
 (1.35,2.14):"4",(2.00,2.14):"497",(3.00,2.14):"479",(4.00,2.14):"14",(4.82,2.14):"2.8%",
 (1.35,2.41):"4",(2.00,2.41):"303",(3.00,2.41):"296",(4.00,2.41):"3",(4.82,2.41):"1.0%",
 (1.35,2.67):"79",(2.00,2.67):"4,736",(3.00,2.67):"4,587",(4.00,2.67):"88",(4.82,2.67):"1.9%",
 (6.80,1.60):"C.잠실롯데월드몰",(8.80,1.60):"197건",
 (6.80,2.07):"S.용산아이파크몰점",(8.80,2.07):"193건",
 (6.80,2.57):"S.롯데창원점",(8.80,2.57):"186건",
 (10.20,1.60):"S.방배점",(12.20,1.60):"8건",
 (10.20,2.07):"I.잠실롯데월드몰점",(12.20,2.07):"6건",
 (10.20,2.57):"H.롯데수원점",(12.20,2.57):"5건",
 (0.35,3.83):"S.방배점",(1.60,3.83):"8건",
 (2.13,3.83):"밥·반찬 양 부족/가격 대비 양 아쉬움/LA갈비 단짠·생선 온도/식기 위생(머리카락)/특별후식(인절미 아이스크림) 사전 안내 부족/홀서빙 응대",
 (8.73,3.83):"밥·반찬 제공량·리필 동선 개선, 고기·생선 품질·온도 점검, 식기 위생 점검, 특별메뉴·주차 사전 안내 강화",
 (0.35,4.26):"I.잠실롯데월드몰점",(1.60,4.26):"6건",
 (2.13,4.26):"스시 신선도·상태 불만(비림)/가격 대비 퀄리티·크기/주방복·금 간 컵 위생/늦은 시간 품절 다수/응대 미흡",
 (8.73,4.26):"원물·스시 신선도 관리 강화, 주방복·식기 위생 점검, 품절 예측·재고 관리, 응대 교육",
 (0.35,4.69):"H.롯데수원점",(1.60,4.69):"5건",
 (2.13,4.69):"가격 대비 양 부족(반복 지적)/오주문 회수 시 응대 당황",
 (8.73,4.69):"세트 구성·제공량 검토, 오주문 처리(제공 음식 회수 지양) 응대 매뉴얼화",
 (0.35,6.09):"음식 맛 3,716",(0.35,6.55):"매장 청결·분위기 1,235",(0.35,7.00):"직원 친절도 823",
 (2.18,6.11):"메뉴 품질·맛 30",(2.18,6.54):"직원 서비스 24",(2.18,7.00):"가격·양 18",
}
used=set()
for shp in list(slide.shapes):
    if not shp.has_text_frame: continue
    try: key=(round(Emu(shp.left).inches,2),round(Emu(shp.top).inches,2))
    except: continue
    if key in M:
        setrun(shp,M[key]); used.add(key)
print("unmatched:",sorted(set(M)-used))

# Section 5 table
tbl=None
for shp in slide.shapes:
    if shp.has_table: tbl=shp.table
rows=[
 ["품질·위생\n점검","메뉴 품질·위생 불만 직접 확인 및 개선","월 1회 정기 점검\n불만 급증 매장 수시 점검","부정 리뷰 건수\n위생 체크리스트 점수"],
 ["서비스\n응대 교육","직원 응대 불만 개선 및 친절도 향상","월별 TOP 불만 유형을 교육 주제로 선정","친절도 리뷰 추이\n서비스 개선 보고서"],
 ["메뉴 구성·\n양 검토","가격 대비 양 불만 대응 및 만족도 개선","메뉴 양·세트 구성 재검토\n특별메뉴 사전 안내 강화","가격·양 부정 리뷰 추이"],
]
def fmt_of(cell):
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            col=None
            try:
                if r.font.color and r.font.color.type is not None: col=r.font.color.rgb
            except: pass
            return (r.font.size, r.font.bold, r.font.name, col)
    return (Pt(9),None,None,None)
for ri,rowdata in enumerate(rows,start=1):
    for ci,val in enumerate(rowdata):
        cell=tbl.cell(ri,ci)
        sz,bold,name,col=fmt_of(cell)
        align=cell.text_frame.paragraphs[0].alignment
        lines=val.split("\n")
        cell.text=lines[0]
        for extra in lines[1:]:
            p=cell.text_frame.add_paragraph(); p.text=extra
        for p in cell.text_frame.paragraphs:
            p.alignment=align
            for r in p.runs:
                if sz: r.font.size=sz
                if bold is not None: r.font.bold=bold
                if name: r.font.name=name
                if col is not None:
                    try: r.font.color.rgb=col
                    except: pass

# Replace keyword image
pic=None
for shp in slide.shapes:
    if shp.shape_type==13 and round(Emu(shp.top).inches,2)>5.0:
        pic=shp
L,T,Wd,Hd=pic.left,pic.top,pic.width,pic.height
pic._element.getparent().remove(pic._element)
slide.shapes.add_picture(WD+r"\kw_new.png",L,T,Wd,Hd)

prs.save(OUT)
print("SAVED:",OUT)
