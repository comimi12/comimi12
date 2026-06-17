# -*- coding: utf-8 -*-
import io
from pptx import Presentation
from pptx.util import Emu
import openpyxl
WD=r"C:\Users\owner\do-better-workspace-v2\00-inbox\voc-may-ppt"
out = io.open(WD+r"\ppt_dump.txt","w",encoding="utf-8")
path = r"C:\Users\owner\Desktop\교육팀\0. 온라인리뷰\★ 26년  VOC 리뷰★\온라인 리뷰_26.05월(05.01~05.31)_V2.pptx"
prs = Presentation(path)
for i, slide in enumerate(prs.slides):
    out.write(f"\n##### SLIDE {i+1} #####\n")
    for shp in slide.shapes:
        try:
            pos=f"({Emu(shp.left).inches:.2f},{Emu(shp.top).inches:.2f}) {Emu(shp.width).inches:.2f}x{Emu(shp.height).inches:.2f}"
        except: pos="?"
        if shp.has_text_frame and shp.text_frame.text.strip():
            fi=""
            for p in shp.text_frame.paragraphs:
                for r in p.runs:
                    sz = r.font.size.pt if r.font.size else "?"
                    col = ""
                    try:
                        if r.font.color and r.font.color.rgb: col=str(r.font.color.rgb)
                    except: pass
                    fi=f" sz={sz} bold={r.font.bold} col={col} font={r.font.name}"
                    break
                if fi: break
            txt = shp.text_frame.text.replace("\n"," | ")
            out.write(f"[{shp.name}] {pos}{fi}\n   TXT: {txt}\n")
        elif shp.has_table:
            out.write(f"[{shp.name}] TABLE {pos}\n")
            for ri,row in enumerate(shp.table.rows):
                cells=[c.text for c in row.cells]
                out.write(f"   R{ri}: {cells}\n")
        else:
            out.write(f"[{shp.name}] {shp.shape_type} {pos}\n")
out.close()
p=r"C:\Users\owner\Desktop\교육팀\0. 온라인리뷰\리뷰집계파일\AI매장리뷰_260611\review_analysis_output\review_analysis_tables_202606111853.xlsx"
wb=openpyxl.load_workbook(p, read_only=True, data_only=True)
o=io.open(WD+r"\xlsx_sheets.txt","w",encoding="utf-8")
for ws in wb.worksheets:
    o.write(f"SHEET: {ws.title}  dims={ws.max_row}x{ws.max_column}\n")
o.close()
print("done")
