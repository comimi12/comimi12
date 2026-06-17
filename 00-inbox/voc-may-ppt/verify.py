# -*- coding: utf-8 -*-
import io
from pptx import Presentation
from pptx.util import Emu
OUT=r"C:\Users\owner\Desktop\교육팀\0. 온라인리뷰\★ 26년  VOC 리뷰★\온라인 리뷰_26.05월(05.01~05.31)_V3.pptx"
prs=Presentation(OUT)
o=io.open(r"C:\Users\owner\do-better-workspace-v2\00-inbox\voc-may-ppt\v3_dump.txt","w",encoding="utf-8")
s=prs.slides[1]
for shp in s.shapes:
    if shp.has_table:
        o.write("[TABLE]\n")
        for r in shp.table.rows:
            o.write("  "+str([c.text.replace(chr(10),"⏎") for c in r.cells])+"\n")
    elif shp.has_text_frame and shp.text_frame.text.strip():
        x=round(Emu(shp.left).inches,2); y=round(Emu(shp.top).inches,2)
        o.write(f"({x},{y}) {shp.text_frame.text}\n")
    elif shp.shape_type==13:
        o.write(f"[PIC] ({round(Emu(shp.left).inches,2)},{round(Emu(shp.top).inches,2)}) {round(Emu(shp.width).inches,2)}x{round(Emu(shp.height).inches,2)}\n")
o.close()
print("ok")
