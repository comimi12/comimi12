# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Emu
path=r"C:\Users\owner\Desktop\교육팀\0. 온라인리뷰\★ 26년  VOC 리뷰★\온라인 리뷰_26.05월(05.01~05.31)_V2.pptx"
prs=Presentation(path)
WD=r"C:\Users\owner\do-better-workspace-v2\00-inbox\voc-may-ppt"
for si,slide in enumerate(prs.slides):
    for shp in slide.shapes:
        if shp.shape_type==13:  # picture
            img=shp.image
            fn=WD+f"\s{si+1}_{shp.shape_id}.{img.ext}"
            open(fn,"wb").write(img.blob)
            print(f"slide{si+1} {shp.name} pos=({Emu(shp.left).inches:.2f},{Emu(shp.top).inches:.2f}) size=({Emu(shp.width).inches:.2f}x{Emu(shp.height).inches:.2f}) -> {fn} ({len(img.blob)}B)")
